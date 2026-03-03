from __future__ import annotations
import os
import pytesseract
from PIL import Image
from pdf2image import convert_from_path
import docx
import openpyxl
from pptx import Presentation
from typing import Optional

import base64
from groq import Groq
from dotenv import load_dotenv

# Load environment variables from backend/.env
env_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'backend', '.env')
load_dotenv(dotenv_path=env_path)

def encode_image(image_path: str) -> str:
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def extract_text_from_image(file_path: str) -> str:
    """Extract text from an image using Groq Vision LLM."""
    api_key = os.getenv("GROQ_API_KEY")
    if not api_key:
        # Fallback to Tesseract if no API key
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image).strip()
        except:
            return ""

    client = Groq(api_key=api_key)
    base64_image = encode_image(file_path)

    try:
        completion = client.chat.completions.create(
            model="llama-3.2-11b-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Extract all text from this business document image. Maintain the structure as much as possible. If it's a French document, keep the French text."},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}",
                            },
                        },
                    ],
                }
            ],
            temperature=0.1,
            max_tokens=2048
        )
        return completion.choices[0].message.content.strip()
    except Exception as e:
        print(f"Groq Vision Error: {e}")
        # Fallback to Tesseract
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image).strip()
        except:
            return ""

import fitz  # PyMuPDF

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF. Attempts native extraction first, then OCR if needed."""
    try:
        # 1. Try native text extraction with PyMuPDF
        doc = fitz.open(file_path)
        native_text = ""
        for page in doc:
            native_text += page.get_text()
        doc.close()
        
        # If we got substantial text, return it
        if len(native_text.strip()) > 50:
            return native_text.strip()
            
        # 2. Fallback to OCR if native text is too short (likely scanned PDF)
        # Note: This requires poppler and tesseract installed on the system
        pages = convert_from_path(file_path)
        ocr_text = ""
        for page in pages:
            ocr_text += pytesseract.image_to_string(page)
        return ocr_text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return ""

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from a .docx file."""
    try:
        doc = docx.Document(file_path)
        text = [para.text for para in doc.paragraphs]
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_xlsx(file_path: str) -> str:
    """Extract text from an .xlsx file."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text.append(" ".join([str(cell) for cell in row if cell is not None]))
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error extracting text from XLSX: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""

def extract_text(file_path: str) -> tuple[str, bool]:
    """General extraction function that routes based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    ocr_used = False
    
    if ext in [".jpg", ".jpeg", ".png", ".tiff"]:
        text = extract_text_from_image(file_path)
        ocr_used = True
    elif ext == ".pdf":
        text = extract_text_from_pdf(file_path)
        ocr_used = True
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        text = extract_text_from_xlsx(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        text = ""
    
    return text, ocr_used
